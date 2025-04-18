import json
from openai import OpenAI
import os
import base64
import os
from pptx import Presentation
import json
import pandas as pd
from openai.types.responses import ResponseTextDeltaEvent
from my_agents import csv_agent
from agents import Runner
import typing as t

client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))


async def extract_text_from_file(file_path, file_type) -> dict[str, t.Any]:
    match file_type:
        case "pdf":
            cache_file_path = file_path + ".cache"
            if os.path.exists(cache_file_path):
                with open(cache_file_path, "r") as f:
                    file_id = f.read().strip()
                print("using cached pdf file id for", file_path)
                return {
                    "type": "input_file",
                    "file_id": file_id,
                }
            print("upload pdf")
            with open(file_path, "rb") as f:
                file = client.files.create(file=f, purpose="user_data")
                file_id = file.id
            with open(cache_file_path, "w") as f:
                f.write(file_id)
            return {
                "type": "input_file",
                "file_id": file_id,
            }
        case "csv":
            cache_file_path = file_path + ".cache"
            if not os.path.exists(cache_file_path):
                df = pd.read_csv(file_path)
                csv_str = df.to_string(index=False)
                print("summarize CSV")
                result = await Runner.run(csv_agent, input=csv_str)
                output = result.final_output
                # async for event in result.stream_events():
                #     if event.type == "raw_response_event" and isinstance(event.data, ResponseTextDeltaEvent):
                #         print(event.data.delta, end="", flush=True)
                #         output += event.data.delta
                print("got summary:", output)
                with open(cache_file_path, "w") as f:
                    f.write(output)
                print("saved to", cache_file_path)
            else:
                with open(cache_file_path) as f:
                    output = f.read()
            json_part = output.split("------")[-1].strip()
            try:
                json_obj = json.loads(json_part)
                json_text = json.dumps(json_obj)
                # TODO: handle this error. Try to fix. Maybe I need to retry the
                # agent. Better would be to fix the agent always respond with
                # correct JSON
            except Exception:
                json_text = json_part
            return {
                "type": "input_text",
                "text": json_text,
            }
        case "pptx":
            prs = Presentation(file_path)
            extracted_text = "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            return {
                "type": "input_text",
                "text": f"the {file_path} presentation content: {extracted_text}",
            }
        case "png" | "jpg":
            return {
                "type": "input_image",
                "image_url": f"data:image/{file_type};base64,{encode_base64(file_path)}",
            }
        case _:
            raise NotImplementedError(f"Not known processing method for {file_type}")


async def extract_data_from_files(directory) -> list[dict]:
    extracted_data = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        file_type = filename.split(".")[-1]
        if file_type == "cache":
            continue
        data = await extract_text_from_file(file_path, file_type)
        if data is not None:
            extracted_data.append(data)
    return extracted_data


# Function to encode the image
def encode_base64(file_path):
    with open(file_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def generate_presentation_story(prompt, extracted_data):

    input_data = [
        {
            "role": "user",
            "content": prompt,
        },
        {"role": "user", "content": extracted_data},
    ]
    print("RUN WITH", input_data)
    response = client.responses.create(
        model="gpt-4.5-preview", 
        instructions="""Create a presentation story based on user prompt. 
        Response with json that should be an array of slides. 
        You can create new slides, generate new content.
        Each slide has template_id - the name of template which is used for the
        slide. Analyze all files and try to use them. User can send json data
        that is actually aggregation info of a csv file. Use that json as a
        source for some total or mean values to inject them in output
        presentation. """,
        input=input_data,
        temperature=0.6
    )
    return response.output_text


def create_json_presentation(story):
    slides = story_to_slides(story)
    templates = load_templates()
    formatted_slides = []

    for slide in slides:
        template = select_best_template(slide, templates)
        formatted_slide = slide  # format_slide_with_template(slide, template)
        formatted_slides.append(formatted_slide)

    return formatted_slides


def load_templates():
    with open("templates.json", "r") as file:
        return json.load(file)


def story_to_slides(story):
    return [slide.strip() for slide in story.split("\n\n") if slide.strip()]


def select_best_template(slide, templates):
    # Implement logic to choose the best template based on content
    return templates[0]


def format_slide_with_template(slide, template):
    # Replace placeholders with actual content
    formatted_slide = template.copy()
    formatted_slide.update({"content": slide})
    return formatted_slide
